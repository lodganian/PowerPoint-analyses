import sys
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE


def take_geometric_shapes(_slide):
    geometric_shapes = []
    for shape in _slide.shapes:
        if not shape.is_placeholder:
            geometric_shapes.append(shape)
    return geometric_shapes


def compare_colors(_sample_color, _tested_color):
    scared = 0
    if _sample_color.type != _tested_color.type:
        return 0

    if _sample_color.type == MSO_COLOR_TYPE.RGB:
        if _sample_color.rgb == _tested_color.rgb:
            scared += 1

    if _sample_color.theme_color == _tested_color.theme_color:
        scared += 1

    if _sample_color.brightness == _tested_color.brightness:
        scared += 1

    return scared


def take_score_from_slide(_sample_shapes, _tested_shapes):
    scored = 0
    for i in range(min(len(_sample_shapes), len(_tested_shapes))):
        scored += compare_two_geometric_shapes(_sample_shapes[i], _tested_shapes[i])
    return scored


def compare_shape_fill(_sample_fill, _tested_fill):
    scored = 0
    if _sample_fill.type == _tested_fill.type:
        scored += 1
    else:
        return 0

    if _sample_fill.type == MSO_FILL_TYPE.GRADIENT:
        if abs(_sample_fill.gradient_angle - _tested_fill.gradient_angle) <= 2:
            scored += 1
        if _sample_fill.gradient_stops == _tested_fill.gradient_stops:
            scored += 1
        scored += compare_colors(_sample_fill.fore_color, _tested_fill.fore_color)
        scored += compare_colors(_sample_fill.back_color, _tested_fill.fore_color)

    elif _sample_fill.type == MSO_FILL_TYPE.PATTERNED:
        if _sample_fill.pattern == _tested_fill.pattern:
            scored += 1
        scored += compare_colors(_sample_fill.fore_color, _tested_fill.fore_color)
        scored += compare_colors(_sample_fill.back_color, _tested_fill.fore_color)

    return scored


def compare_shape_line(_sample_line, _tested_line):
    scored = 0
    if _sample_line.dash_style == _tested_line.dash_style:
        scored += 1

    scored += compare_colors(_sample_line.color, _tested_line.color)

    return scored


def compare_shape_offsets(_sample_shape, _tested_shape):
    scored = 0
    left_sample = _sample_shape.left.pt
    top_sample = _sample_shape.top.pt

    left_tested = _tested_shape.left.pt
    top_tested = _tested_shape.top.pt

    if abs(left_sample - left_tested) <= 15:
        scored += 1
    if abs(top_sample - top_tested) <= 15:
        scored += 1

    return scored


def compare_two_geometric_shapes(shape1, shape2):
    scored = 0

    if shape1.shape_type == shape2.shape_type:
        scored += 1

    scored += compare_shape_fill(shape1.fill, shape2.fill)
    scored += compare_shape_line(shape1.line, shape2.line)
    scored += compare_shape_offsets(shape1, shape2)

    return scored


def compare_slide(_sample_slide, _tested_slide):
    sample_shapes = take_geometric_shapes(_sample_slide)
    tested_shapes = take_geometric_shapes(_tested_slide)

    return take_score_from_slide(sample_shapes, tested_shapes)


def analysis_slides(_sample_slides, _tested_slides):
    scored = 0
    for i in range(min(len(_sample_slides), len(_tested_slides))):
        scored += compare_slide(_sample_slides[i], _tested_slides[i])
    return scored


def compare_shadow():
    return 0


sample = Presentation(sys.argv[1])
tested = Presentation(sys.argv[2])

sample_slides = sample.slides
tested_slides = tested.slides


sys.exit(analysis_slides(sample_slides, tested_slides))
